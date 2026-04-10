import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Briefing Studio–inspired palette (slate + teal)
_BG = RGBColor(15, 23, 42)
_CARD = RGBColor(30, 41, 59)
_CARD_BORDER = RGBColor(51, 65, 85)
_ACCENT = RGBColor(45, 212, 191)
_TEXT = RGBColor(248, 250, 252)
_MUTED = RGBColor(148, 163, 184)
_WHITE = RGBColor(255, 255, 255)


def _full_bleed(slide, prs: Presentation, rgb: RGBColor) -> None:
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        int(prs.slide_width),
        int(prs.slide_height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb
    shp.line.fill.background()


def _footer_teal_line(slide, prs: Presentation) -> None:
    h = int(Inches(0.06))
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        int(prs.slide_height) - h,
        int(prs.slide_width),
        h,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _ACCENT
    bar.line.fill.background()


def _content_card(slide, prs: Presentation) -> Any:
    m_l, m_t = Inches(0.55), Inches(0.5)
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(m_l),
        int(m_t),
        int(prs.slide_width) - 2 * int(m_l),
        int(prs.slide_height) - int(m_t) - int(Inches(0.45)),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = _CARD
    card.line.color.rgb = _CARD_BORDER
    return card


def _heading_box(slide, y_inches: float, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.85), Inches(y_inches), Inches(11.5), Inches(0.75))
    tf = tb.text_frame
    tf.margin_bottom = tf.margin_top = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = _ACCENT
    p.alignment = PP_ALIGN.LEFT
    # Thin accent underline via small shape
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(Inches(0.85)),
        int(Inches(y_inches + 0.72)),
        int(Inches(2.4)),
        int(Inches(0.04)),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _ACCENT
    line.line.fill.background()


def _body_box(slide, top: float, height: float) -> Any:
    return slide.shapes.add_textbox(Inches(0.95), Inches(top), Inches(11.4), Inches(height))


def build_deck(summary: dict[str, Any], title: str, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Title
    s0 = prs.slides.add_slide(blank)
    _full_bleed(s0, prs, _BG)
    _footer_teal_line(s0, prs)
    accent = s0.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        int(Inches(0)),
        int(Inches(0.14)),
        int(prs.slide_height),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _ACCENT
    accent.line.fill.background()

    tbox = s0.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(11.2), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "Meeting briefing"
    p0.font.size = Pt(46)
    p0.font.bold = True
    p0.font.color.rgb = _WHITE
    p1 = tf.add_paragraph()
    p1.text = title
    p1.font.size = Pt(22)
    p1.font.color.rgb = _MUTED
    p1.space_before = Pt(12)
    p2 = tf.add_paragraph()
    p2.text = "Confidential — internal use"
    p2.font.size = Pt(13)
    p2.font.color.rgb = _ACCENT
    p2.space_before = Pt(28)

    # Executive summary
    s1 = prs.slides.add_slide(blank)
    _full_bleed(s1, prs, _BG)
    _footer_teal_line(s1, prs)
    _content_card(s1, prs)
    _heading_box(s1, 0.68, "Executive summary")
    body = _body_box(s1, 1.55, 5.35)
    btf = body.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.text = str(summary.get("executive_summary") or "").strip() or "—"
    bp.font.size = Pt(18)
    bp.font.color.rgb = _TEXT
    bp.line_spacing = 1.2

    def add_bullets(heading: str, bullets: list[str]) -> None:
        sl = prs.slides.add_slide(blank)
        _full_bleed(sl, prs, _BG)
        _footer_teal_line(sl, prs)
        _content_card(sl, prs)
        _heading_box(sl, 0.68, heading)
        box = _body_box(sl, 1.45, 5.45)
        btf2 = box.text_frame
        btf2.word_wrap = True
        lines = bullets if bullets else ["—"]
        for i, line in enumerate(lines[:12]):
            para = btf2.paragraphs[0] if i == 0 else btf2.add_paragraph()
            para.text = line
            para.level = 0
            para.font.size = Pt(17)
            para.font.color.rgb = _TEXT
            para.space_after = Pt(10)
            para.line_spacing = 1.15

    add_bullets("Strategic objectives", list(summary.get("objectives", [])[:3]))
    add_bullets("Actionable items", list(summary.get("actionable_items", [])[:3]))
    add_bullets("Next steps", list(summary.get("next_steps", [])))

    # Human review
    s_h = prs.slides.add_slide(blank)
    _full_bleed(s_h, prs, _BG)
    _footer_teal_line(s_h, prs)
    _content_card(s_h, prs)
    _heading_box(s_h, 0.68, "Human review")
    hb = _body_box(s_h, 1.55, 5.35)
    htf = hb.text_frame
    htf.word_wrap = True
    hp = htf.paragraphs[0]
    hp.text = str(summary.get("human_review_notes") or "").strip() or "—"
    hp.font.size = Pt(17)
    hp.font.color.rgb = _MUTED
    hp.line_spacing = 1.2

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def write_summary_json(
    summary: dict[str, Any],
    path: Path,
    *,
    meta: dict[str, Any] | None = None,
) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    out = dict(summary)
    if meta:
        out["_briefing_meta"] = meta
    path.write_text(json.dumps(out, indent=2, ensure_ascii=False), encoding="utf-8")
